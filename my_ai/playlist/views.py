# from django.shortcuts import render
# import subprocess
# import json
# from rest_framework.views import APIView
# from rest_framework.response import Response
# from rest_framework import status


# class ExtractPlaylistView(APIView):
#     def post(self, request):
#         playlist_url = request.data.get("playlistUrl")

#         if not playlist_url:
#             return Response(
#                 {"error": "No playlist URL provided"},
#                 status=status.HTTP_400_BAD_REQUEST,
#             )

#         try:
#             # Run yt-dlp to extract playlist metadata
#             command = ["yt-dlp", "-j", "--flat-playlist", playlist_url]
#             result = subprocess.run(command, capture_output=True, text=True, check=True)

#             videos = []
#             for line in result.stdout.strip().split("\n"):
#                 data = json.loads(line)
#                 videos.append(
#                     {
#                         "title": data.get("title"),
#                         "url": f"https://www.youtube.com/watch?v={data.get('id')}",
#                     }
#                 )
#             print(videos)
#             return Response({"videos": videos}, status=status.HTTP_200_OK)

#         except subprocess.CalledProcessError as e:
#             return Response(
#                 {"error": e.stderr}, status=status.HTTP_500_INTERNAL_SERVER_ERROR
#             )


# # views.py
# # views.py
# from django.http import JsonResponse
# from rest_framework.views import APIView
# import json, math, re
# from collections import Counter
# from urllib.parse import urlparse, parse_qs

# # For PDF
# import requests
# from io import BytesIO
# from PyPDF2 import PdfReader

# # Load difficulty JSON once
# import json
# import re
# import subprocess
# from collections import Counter
# from urllib.parse import urlparse, parse_qs
# from typing import List, Tuple, Optional
# import PyPDF2
# from pptx import Presentation
# import requests
# import tempfile
# import os
# from sklearn.feature_extraction.text import ENGLISH_STOP_WORDS
# import re
# from typing import List
# from typing import List
# from sklearn.feature_extraction.text import TfidfVectorizer
# import numpy as np
# import re
# import nltk


# def safe_nltk_download(resource):
#     try:
#         nltk.data.find(resource)
#     except LookupError:
#         nltk.download(resource.split("/")[-1], quiet=True)


# # Ensure all NLTK models are available
# for res in [
#     "tokenizers/punkt",
#     "tokenizers/punkt_tab",
#     "taggers/averaged_perceptron_tagger",
#     "taggers/averaged_perceptron_tagger_eng",
# ]:
#     safe_nltk_download(res)
# # Load difficulty map once
# with open("data/difficulty_map.json", "r") as f:
#     DIFFICULTY_MAP = json.load(f)
# import joblib


# def clean_youtube_url(url: str) -> str:
#     """Normalize YouTube URLs to standard watch?v=VIDEOID format."""
#     try:
#         parsed = urlparse(url)
#         hostname = parsed.hostname.lower() if parsed.hostname else ""

#         if hostname in ("youtu.be", "www.youtu.be"):
#             video_id = parsed.path.lstrip("/")
#             return f"https://www.youtube.com/watch?v={video_id}"

#         if "youtube.com" in hostname:
#             v = parse_qs(parsed.query).get("v")
#             if v:
#                 return f"https://www.youtube.com/watch?v={v[0]}"
#             match = re.search(r"/(shorts|embed|live)/([a-zA-Z0-9_-]{11})", parsed.path)
#             if match:
#                 return f"https://www.youtube.com/watch?v={match.group(2)}"
#     except Exception as e:
#         print(f"[ERROR] Invalid YouTube URL: {url} -> {e}")
#     return url


# class TaskDifficultyCalculator:

#     def __init__(self, difficulty_map: List[dict]):
#         self.difficulty_map = difficulty_map
#         self.model = joblib.load("difficulty_model.pkl")

#     def extract_keywords_from_text(
#         self, text: str, max_keywords: int = 30
#     ) -> List[str]:
#         if not isinstance(text, str) or not text.strip():
#             return []

#         # Clean the text
#         text = re.sub(r"http\S+|www\S+|youtu\.be\S+", "", text)
#         text = re.sub(r"[^a-zA-Z\s]", " ", text).lower().strip()
#         if not text:
#             return []

#         # Tokenize and extract relevant words
#         words = nltk.word_tokenize(text)
#         tagged = nltk.pos_tag(words)
#         topic_words = [
#             w for w, pos in tagged if pos.startswith("NN") or pos.startswith("JJ")
#         ]
#         if not topic_words:
#             return []

#         filtered_text = " ".join(topic_words)
#         if not filtered_text.strip():
#             return list(set(topic_words))[:max_keywords]

#         try:
#             vectorizer = TfidfVectorizer(stop_words="english", max_features=1000)
#             tfidf_matrix = vectorizer.fit_transform([filtered_text])

#             if tfidf_matrix.shape[1] == 0:
#                 # Fallback: return top unique frequent words
#                 freq = Counter(topic_words)
#                 return [w for w, _ in freq.most_common(max_keywords)]

#             feature_names = np.array(vectorizer.get_feature_names_out())
#             scores = tfidf_matrix.toarray()[0]
#             top_indices = np.argsort(scores)[::-1][:100]
#             print(feature_names[top_indices].tolist())
#             return feature_names[top_indices].tolist()

#         except Exception as e:
#             print(f"[ERROR] extract_keywords_from_text: {e}")
#             return []

#     def get_youtube_keywords(self, url: str) -> Tuple[List[str], float]:
#         """Extract keywords and duration from YouTube using yt-dlp."""
#         url = clean_youtube_url(url)
#         try:
#             command = ["yt-dlp", "-j", "--skip-download", url]
#             result = subprocess.run(command, capture_output=True, text=True, check=True)
#             data = json.loads(result.stdout)
#             title = data.get("title", "")
#             description = data.get("description", "")
#             duration_hours = data.get("duration", 0) / 3600
#             keywords = self.extract_keywords_from_text(f"{title} {description}")
#             return keywords, duration_hours
#         except Exception as e:
#             print(f"[ERROR] get_youtube_keywords: {e}")
#             return [], 0.1

#     def get_document_keywords(self, file_path: str) -> Tuple[List[str], float]:
#         """
#         Extract keywords and estimated reading hours from a local PDF/PPTX file.
#         - PDFs: use number of words in all pages
#         - PPTX: use number of words in all slides
#         - Returns keywords and estimated hours for an average reader
#         """
#         text = ""
#         try:
#             ext = os.path.splitext(file_path)[1].lower()
#             if ext == ".pdf":
#                 with open(file_path, "rb") as f:
#                     reader = PyPDF2.PdfReader(f)
#                     text = " ".join(page.extract_text() or "" for page in reader.pages)
#             elif ext == ".pptx":
#                 prs = Presentation(file_path)
#                 text = " ".join(
#                     shape.text
#                     for slide in prs.slides
#                     for shape in slide.shapes
#                     if hasattr(shape, "text")
#                 )
#             elif ext == ".ppt":
#                 print("[WARN] PPT files not directly supported. Convert to PPTX first.")
#         except Exception as e:
#             print(f"[ERROR] get_document_keywords: {e}")

#         # Extract words
#         words = re.findall(r"\b\w+\b", text)
#         num_words = len(words)

#         # Estimate hours based on average reading speed (~200 wpm)
#         minutes = num_words / 20
#         hours = max(minutes / 60, 0.1)  # minimum 0.1 hour

#         # Extract keywords
#         keywords = self.extract_keywords_from_text(text)
#         return keywords, hours

#     def get_document_keywords_from_url(self, url: str) -> Tuple[List[str], float]:
#         """
#         Download a public PDF/PPTX/PPT file, save to temp file, extract keywords and estimated hours.
#         Works on Windows without Permission errors.
#         """
#         try:
#             response = requests.get(url, stream=True)
#             response.raise_for_status()
#             ext = os.path.splitext(urlparse(url).path)[1].lower()
#             if ext not in (".pdf", ".pptx", ".ppt"):
#                 print(f"[WARN] Unsupported document type: {url}")
#                 return [], 0.1

#             # Create temp file safely on Windows
#             with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
#                 for chunk in response.iter_content(chunk_size=8192):
#                     tmp_file.write(chunk)
#                 tmp_path = tmp_file.name  # Save path for later

#             # Now safely open the temp file in your document handler
#             keywords, hours = self.get_document_keywords(tmp_path)

#             # Clean up temp file
#             os.remove(tmp_path)
#             return keywords, hours

#         except Exception as e:
#             print(f"[ERROR] Failed to process document URL {url}: {e}")
#             return [], 0.1

#     def compute_difficulty(self, video_data: List[dict]) -> float:
#         """Compute normalized difficulty score from keywords and hours."""
#         all_keywords = []
#         total_hours = 0
#         for v in video_data:
#             all_keywords.extend([kw.lower() for kw in v.get("keywords", [])])
#             total_hours += v.get("hours", 0)
#         predictions = self.model.predict(all_keywords)
#         print(predictions)
#         base_difficulty = round(predictions.mean(), 2)
#         print(base_difficulty)
#         raw_score = base_difficulty * total_hours
#         normalized_score = min((raw_score / (40 * 10)) * 100, 100)
#         return round(normalized_score, 1)

#     def compute_document_difficulty(
#         self, keywords: List[str], hours_estimate: float
#     ) -> float:
#         """
#         Compute difficulty for a document based on extracted keywords and estimated reading hours.

#         Args:
#            keywords (List[str]): extracted keywords from the PDF/PPTX/PPT
#            hours_estimate (float): estimated reading hours based on word count/pages

#         Returns:
#            float: normalized difficulty score (0-100)
#         """
#         # Step 1: Compute base difficulty from keywords
#         scores = []
#         for entry in self.difficulty_map:
#             entry_keywords = [kw.lower() for kw in entry.get("keywords", [])]
#             for kw in keywords:
#                 if kw in entry_keywords or any(
#                     kw in ek or ek in kw for ek in entry_keywords
#                 ):
#                     scores.append(entry["difficulty"])
#                     break

#         base_difficulty = 5 if not scores else sum(scores) / len(scores)

#         # Step 2: Scale effective hours
#         effective_hours = hours_estimate * (
#             base_difficulty / 5
#         )  # adjust hours based on difficulty

#         # Step 3: Normalize score to 0-100
#         raw_score = base_difficulty * effective_hours
#         normalized_score = min((raw_score / (40 * 10)) * 100, 100)

#         print(
#             f"[DEBUG] Document: Keywords={len(keywords)}, Hours={hours_estimate:.2f}, "
#             f"BaseDifficulty={base_difficulty:.2f}, Score={normalized_score:.1f}"
#         )

#         return round(normalized_score, 1)


#     def calculate(self, source: str) -> Tuple[List[str], Optional[float]]:
#      """Calculate difficulty for YouTube URL or public document URL."""
#      try:
#         parsed = urlparse(source)

#         if parsed.scheme in ("http", "https"):
#             # --- Handle YouTube Links ---
#             if "youtube.com" in (parsed.hostname or "") or "youtu.be" in (
#                 parsed.hostname or ""
#             ):
#                 keywords, hours = self.get_youtube_keywords(source)
#                 difficulty = self.compute_difficulty(
#                     [{"keywords": keywords, "hours": hours}]
#                 )
#                 return keywords, difficulty

#             # --- Handle Document Links ---
#             elif source.lower().endswith((".pdf", ".pptx", ".ppt")):
#                 keywords, hours = self.get_document_keywords_from_url(source)
#                 difficulty = self.compute_document_difficulty(keywords, hours)
#                 return keywords, difficulty

#         # --- Fallback for unsupported sources ---
#         print(f"[DEBUG] Unsupported source: {source}")
#         return [], None

#      except Exception as e:
#         print(f"[ERROR] Exception in calculate(): {e}")
#         return [], None


# class CalculateTaskDifficultyView(APIView):
#     def post(self, request):
#         url = request.data.get("url")
#         print(f"[DEBUG] Received URL: {url}")
#         if not url:
#             return JsonResponse({"error": "URL not provided"}, status=400)

#         calculator = TaskDifficultyCalculator(DIFFICULTY_MAP)
#         try:
#             keywords, difficulty = calculator.calculate(url)
#             if difficulty is None:
#                 print(f"[DEBUG] Unsupported URL type")
#                 return JsonResponse({"error": "Unsupported URL"}, status=400)
#             print(f"[DEBUG] Returning difficulty: {difficulty}")
#             return JsonResponse({"keywords": keywords, "difficulty": difficulty})
#         except Exception as e:
#             print(f"[ERROR] Exception during calculation: {e}")
#             return JsonResponse({"error": str(e)}, status=500)
from django.http import JsonResponse
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status

import subprocess
import json
import os
import re
import tempfile
from collections import Counter
from urllib.parse import urlparse, parse_qs
from typing import List, Tuple, Optional

import requests
import numpy as np
import joblib
import nltk
from sklearn.feature_extraction.text import TfidfVectorizer
from PyPDF2 import PdfReader
from pptx import Presentation

# =====================================================
# 🔹 ONE-TIME GLOBAL INITIALIZATION (CRITICAL)
# =====================================================

# Ensure required NLTK assets ONCE
import threading

_nltk_lock = threading.Lock()

REQUIRED_NLTK = [
    "tokenizers/punkt",
    "taggers/averaged_perceptron_tagger",
]

with _nltk_lock:
    for res in REQUIRED_NLTK:
        try:
            nltk.data.find(res)
        except LookupError:
            nltk.download(res.split("/")[-1], quiet=True)

# -------------------------------
# 🔹 ML MODEL (Google Drive)
# -------------------------------

MODEL_PATH = "difficulty_model.pkl"

MODEL_URL = os.getenv(
    "MODEL_URL",
    "https://drive.google.com/uc?export=download&id=1aKU-7u94cw-pJo5GF57STRZlV0HRN3Rm"
)

def download_model_once():
    if os.path.exists(MODEL_PATH):
        return

    print("⬇️ Downloading ML model from Google Drive...")
    response = requests.get(MODEL_URL, stream=True, timeout=120)
    response.raise_for_status()

    content_type = response.headers.get("Content-Type", "")
    if "text/html" in content_type:
        raise RuntimeError(
            "Google Drive returned HTML instead of model file. "
            "Check sharing permissions."
        )

    with open(MODEL_PATH, "wb") as f:
        for chunk in response.iter_content(chunk_size=8192):
            if chunk:
                f.write(chunk)

    print("✅ ML model downloaded")

download_model_once()

MODEL = joblib.load(MODEL_PATH)

# -------------------------------
# 🔹 DIFFICULTY MAP
# -------------------------------

with open("data/difficulty_map.json", "r") as f:
    DIFFICULTY_MAP = json.load(f)

# -------------------------------
# 🔹 SIMPLE IN-MEMORY CACHE
# -------------------------------

CACHE = {}
CACHE_LIMIT = 500

def cache_get(key):
    return CACHE.get(key)

def cache_set(key, value):
    if len(CACHE) >= CACHE_LIMIT:
        CACHE.pop(next(iter(CACHE)))
    CACHE[key] = value
# =====================================================
# 🔹 PLAYLIST EXTRACTION (UNCHANGED BEHAVIOR)
# =====================================================

class ExtractPlaylistView(APIView):
    def post(self, request):
        playlist_url = request.data.get("playlistUrl")

        if not playlist_url:
            return Response(
                {"error": "No playlist URL provided"},
                status=status.HTTP_400_BAD_REQUEST,
            )

        try:
            command = ["yt-dlp", "-j", "--flat-playlist", playlist_url]
            result = subprocess.run(
                command,
                capture_output=True,
                text=True,
                timeout=30,
                check=True,
            )

            videos = []
            for line in result.stdout.strip().split("\n"):
                data = json.loads(line)
                videos.append(
                    {
                        "title": data.get("title"),
                        "url": f"https://www.youtube.com/watch?v={data.get('id')}",
                    }
                )

            return Response({"videos": videos}, status=status.HTTP_200_OK)

        except subprocess.CalledProcessError as e:
            return Response(
                {"error": e.stderr},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR,
            )


# =====================================================
# 🔹 UTILITIES (UNCHANGED, OPTIMIZED)
# =====================================================

def clean_youtube_url(url: str) -> str:
    try:
        parsed = urlparse(url)
        hostname = parsed.hostname.lower() if parsed.hostname else ""

        if hostname in ("youtu.be", "www.youtu.be"):
            return f"https://www.youtube.com/watch?v={parsed.path.lstrip('/')}"

        if "youtube.com" in hostname:
            v = parse_qs(parsed.query).get("v")
            if v:
                return f"https://www.youtube.com/watch?v={v[0]}"
            match = re.search(r"/(shorts|embed|live)/([a-zA-Z0-9_-]{11})", parsed.path)
            if match:
                return f"https://www.youtube.com/watch?v={match.group(2)}"

    except Exception:
        pass

    return url


# =====================================================
# 🔹 CORE DIFFICULTY CALCULATOR (FULL FEATURE SET)
# =====================================================

class TaskDifficultyCalculator:
    def __init__(self, difficulty_map: List[dict]):
        self.difficulty_map = difficulty_map
        self.model = MODEL  # preloaded model

    def extract_keywords_from_text(
        self, text: str, max_keywords: int = 30
    ) -> List[str]:
        if not isinstance(text, str) or not text.strip():
            return []

        text = re.sub(r"http\S+|[^a-zA-Z\s]", " ", text).lower()
        words = nltk.word_tokenize(text)
        tagged = nltk.pos_tag(words)

        topic_words = [
            w for w, pos in tagged if pos.startswith("NN") or pos.startswith("JJ")
        ]

        if not topic_words:
            return []

        try:
            vectorizer = TfidfVectorizer(stop_words="english", max_features=100)
            tfidf = vectorizer.fit_transform([" ".join(topic_words)])
            return vectorizer.get_feature_names_out().tolist()[:max_keywords]
        except Exception:
            freq = Counter(topic_words)
            return [w for w, _ in freq.most_common(max_keywords)]

    def get_youtube_keywords(self, url: str) -> Tuple[List[str], float]:
        url = clean_youtube_url(url)

        cached = cache_get(url)
        if cached:
            return cached

        try:
            result = subprocess.run(
                ["yt-dlp", "-j", "--skip-download", url],
                capture_output=True,
                text=True,
                timeout=20,
                check=True,
            )
            data = json.loads(result.stdout)

            title = data.get("title", "")
            description = data.get("description", "")
            duration_hours = max(data.get("duration", 600) / 3600, 0.1)

            keywords = self.extract_keywords_from_text(
                f"{title} {description}"
            )

            cache_set(url, (keywords, duration_hours))
            return keywords, duration_hours

        except Exception:
            return [], 0.1

    def get_document_keywords(self, file_path: str) -> Tuple[List[str], float]:
        text = ""
        try:
            ext = os.path.splitext(file_path)[1].lower()

            if ext == ".pdf":
                with open(file_path, "rb") as f:
                    reader = PdfReader(f)
                    text = " ".join(p.extract_text() or "" for p in reader.pages)

            elif ext == ".pptx":
                prs = Presentation(file_path)
                text = " ".join(
                    shape.text
                    for slide in prs.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                )
        except Exception:
            pass

        words = re.findall(r"\w+", text)
        hours = max(len(words) / 200 / 60, 0.1)
        return self.extract_keywords_from_text(text), hours

    def get_document_keywords_from_url(self, url: str) -> Tuple[List[str], float]:
        try:
            response = requests.get(url, timeout=10)
            response.raise_for_status()

            ext = os.path.splitext(urlparse(url).path)[1].lower()
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(response.content)
                tmp_path = tmp.name

            keywords, hours = self.get_document_keywords(tmp_path)
            os.remove(tmp_path)
            return keywords, hours

        except Exception:
            return [], 0.1

    def compute_difficulty(self, video_data: List[dict]) -> float:
        all_keywords = []
        total_hours = 0

        for v in video_data:
            all_keywords.extend(v.get("keywords", []))
            total_hours += v.get("hours", 0)

        if not all_keywords:
            return 5.0

        preds = self.model.predict(all_keywords)
        base = preds.mean()

        raw_score = base * total_hours
        normalized = min((raw_score / 400) * 100, 100)

        return round(normalized, 1)

    def compute_document_difficulty(
        self, keywords: List[str], hours_estimate: float
    ) -> float:
        scores = []
        for entry in self.difficulty_map:
            ek = [k.lower() for k in entry.get("keywords", [])]
            if any(k in ek for k in keywords):
                scores.append(entry["difficulty"])

        base = 5 if not scores else sum(scores) / len(scores)
        effective_hours = hours_estimate * (base / 5)

        raw = base * effective_hours
        normalized = min((raw / 400) * 100, 100)

        return round(normalized, 1)

    def calculate(self, source: str) -> Tuple[List[str], Optional[float]]:
        parsed = urlparse(source)

        if parsed.scheme in ("http", "https"):
            if "youtube.com" in source or "youtu.be" in source:
                keywords, hours = self.get_youtube_keywords(source)
                difficulty = self.compute_difficulty(
                    [{"keywords": keywords, "hours": hours}]
                )
                return keywords, difficulty

            if source.lower().endswith((".pdf", ".pptx", ".ppt")):
                keywords, hours = self.get_document_keywords_from_url(source)
                difficulty = self.compute_document_difficulty(keywords, hours)
                return keywords, difficulty

        return [], None


# =====================================================
# 🔹 API VIEW (UNCHANGED INTERFACE)
# =====================================================

class CalculateTaskDifficultyView(APIView):
    def post(self, request):
        url = request.data.get("url")

        if not url:
            return JsonResponse({"error": "URL not provided"}, status=400)

        cached = cache_get(url)
        if cached:
            keywords, difficulty = cached
            return JsonResponse({"keywords": keywords, "difficulty": difficulty})

        calculator = TaskDifficultyCalculator(DIFFICULTY_MAP)
        keywords, difficulty = calculator.calculate(url)

        if difficulty is None:
            return JsonResponse({"error": "Unsupported URL"}, status=400)

        cache_set(url, (keywords, difficulty))
        return JsonResponse({"keywords": keywords, "difficulty": difficulty})
